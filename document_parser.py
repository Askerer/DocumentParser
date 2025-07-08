import pandas as pd
import openpyxl
from pptx import Presentation
import xml.etree.ElementTree as ET
from PIL import Image
import base64
import io
import os
import re
from datetime import datetime
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument

class DocumentParser:
    def __init__(self):
        self.images_dir = "extracted_images"
        os.makedirs(self.images_dir, exist_ok=True)
    
    def extract_content(self, file_path: str, file_extension: str) -> dict:
        """根據檔案類型擷取內容"""
        try:
            if file_extension in ['.xlsx', '.xls']:
                return self._extract_excel_content(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                return self._extract_powerpoint_content(file_path)
            elif file_extension == '.xml':
                return self._extract_xml_content(file_path)
            elif file_extension == '.pdf':
                return self._extract_pdf_content(file_path)
            elif file_extension in ['.txt', '.md', '.csv', '.docx']:
                return self._extract_text_content(file_path, file_extension)
            else:
                raise ValueError(f"不支援的檔案格式: {file_extension}")
        except Exception as e:
            raise Exception(f"解析檔案時發生錯誤: {str(e)}")
    
    def _extract_excel_content(self, file_path: str) -> dict:
        """擷取Excel檔案內容"""
        try:
            # 使用pandas讀取Excel檔案
            excel_file = pd.ExcelFile(file_path)
            sheets_data = {}
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheets_data[sheet_name] = {
                    'data': df.to_dict('records'),
                    'columns': df.columns.tolist(),
                    'shape': df.shape,
                    'summary': {
                        'total_rows': len(df),
                        'total_columns': len(df.columns),
                        'non_null_counts': df.count().to_dict()
                    }
                }
            
            return {
                'file_type': 'excel',
                'total_sheets': len(excel_file.sheet_names),
                'sheet_names': excel_file.sheet_names,
                'sheets': sheets_data,
                'extraction_time': datetime.now().isoformat()
            }
        except Exception as e:
            raise Exception(f"Excel檔案解析失敗: {str(e)}")
    
    def _extract_powerpoint_content(self, file_path: str) -> dict:
        """擷取PowerPoint檔案內容"""
        try:
            prs = Presentation(file_path)
            slides_data = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = {
                    'slide_number': i + 1,
                    'text_content': [],
                    'shapes_count': len(slide.shapes)
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content['text_content'].append({
                            'type': 'text',
                            'content': shape.text.strip()
                        })
                    elif shape.shape_type == 13:  # 圖片
                        slide_content['text_content'].append({
                            'type': 'image',
                            'description': '圖片內容'
                        })
                
                slides_data.append(slide_content)
            
            return {
                'file_type': 'powerpoint',
                'total_slides': len(prs.slides),
                'slides': slides_data,
                'extraction_time': datetime.now().isoformat()
            }
        except Exception as e:
            raise Exception(f"PowerPoint檔案解析失敗: {str(e)}")
    
    def _extract_xml_content(self, file_path: str) -> dict:
        """擷取XML檔案內容，包含圖片處理"""
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            # 擷取所有文字內容
            text_content = []
            images_extracted = []
            
            def extract_element_content(element, path=""):
                current_path = f"{path}/{element.tag}" if path else element.tag
                
                # 處理文字內容
                if element.text and element.text.strip():
                    text_content.append({
                        'tag': element.tag,
                        'path': current_path,
                        'content': element.text.strip()
                    })
                
                # 處理屬性
                if element.attrib:
                    for attr_name, attr_value in element.attrib.items():
                        text_content.append({
                            'tag': f"{element.tag}@{attr_name}",
                            'path': current_path,
                            'content': attr_value
                        })
                
                # 處理子元素
                for child in element:
                    extract_element_content(child, current_path)
                
                # 處理圖片編碼
                if element.text and self._is_base64_image(element.text):
                    try:
                        image_data = base64.b64decode(element.text)
                        image = Image.open(io.BytesIO(image_data))
                        
                        # 儲存圖片
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        image_filename = f"xml_image_{timestamp}_{len(images_extracted)}.png"
                        image_path = os.path.join(self.images_dir, image_filename)
                        image.save(image_path)
                        
                        images_extracted.append({
                            'tag': element.tag,
                            'path': current_path,
                            'image_path': image_path,
                            'image_size': image.size,
                            'image_format': image.format
                        })
                        
                        text_content.append({
                            'tag': element.tag,
                            'path': current_path,
                            'content': f"[圖片已擷取: {image_filename}]",
                            'image_info': {
                                'path': image_path,
                                'size': image.size,
                                'format': image.format
                            }
                        })
                    except Exception as e:
                        text_content.append({
                            'tag': element.tag,
                            'path': current_path,
                            'content': f"[圖片解析失敗: {str(e)}]"
                        })
            
            extract_element_content(root)
            
            return {
                'file_type': 'xml',
                'root_tag': root.tag,
                'text_content': text_content,
                'images_extracted': images_extracted,
                'total_elements': len(text_content),
                'total_images': len(images_extracted),
                'extraction_time': datetime.now().isoformat()
            }
        except Exception as e:
            raise Exception(f"XML檔案解析失敗: {str(e)}")
    
    def _is_base64_image(self, text: str) -> bool:
        """檢查文字是否為base64編碼的圖片"""
        # 檢查是否為base64編碼
        if not text or len(text) < 100:  # 圖片編碼通常很長
            return False
        
        try:
            # 檢查是否為有效的base64
            base64.b64decode(text)
            # 檢查是否包含圖片特徵
            return any(img_type in text[:50].lower() for img_type in ['data:image', 'png', 'jpg', 'jpeg', 'gif', 'bmp'])
        except:
            return False
    
    def _extract_pdf_content(self, file_path: str) -> dict:
        """擷取PDF檔案內容"""
        try:
            # 使用PyMuPDF (fitz) 進行更準確的PDF解析
            doc = fitz.open(file_path)
            pages_data = []
            total_text = ""
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # 擷取文字
                text = page.get_text()
                total_text += text + "\n"
                
                # 擷取圖片
                images = page.get_images()
                page_images = []
                
                for img_index, img in enumerate(images):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        if pix.n - pix.alpha < 4:  # 檢查是否為RGB或CMYK
                            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                            image_filename = f"pdf_image_p{page_num+1}_{img_index}_{timestamp}.png"
                            image_path = os.path.join(self.images_dir, image_filename)
                            pix.save(image_path)
                            
                            page_images.append({
                                'image_path': image_path,
                                'size': (pix.width, pix.height),
                                'colorspace': pix.colorspace.name
                            })
                        
                        pix = None  # 釋放記憶體
                    except Exception as e:
                        print(f"圖片擷取失敗: {str(e)}")
                
                pages_data.append({
                    'page_number': page_num + 1,
                    'text_content': text.strip(),
                    'images': page_images,
                    'text_length': len(text)
                })
            
            doc.close()
            
            return {
                'file_type': 'pdf',
                'total_pages': len(doc),
                'pages': pages_data,
                'total_text': total_text,
                'total_text_length': len(total_text),
                'extraction_time': datetime.now().isoformat()
            }
        except Exception as e:
            raise Exception(f"PDF檔案解析失敗: {str(e)}")
    
    def _extract_text_content(self, file_path: str, file_extension: str) -> dict:
        """擷取文字檔案內容"""
        try:
            if file_extension == '.docx':
                # 處理Word文件
                doc = DocxDocument(file_path)
                paragraphs = []
                tables_data = []
                
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append({
                            'text': para.text.strip(),
                            'style': para.style.name
                        })
                
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    tables_data.append(table_data)
                
                return {
                    'file_type': 'docx',
                    'paragraphs': paragraphs,
                    'tables': tables_data,
                    'total_paragraphs': len(paragraphs),
                    'total_tables': len(tables_data),
                    'extraction_time': datetime.now().isoformat()
                }
            else:
                # 處理一般文字檔案
                encoding = self._detect_encoding(file_path)
                
                with open(file_path, 'r', encoding=encoding) as file:
                    content = file.read()
                
                lines = content.split('\n')
                
                return {
                    'file_type': file_extension[1:],  # 移除點號
                    'content': content,
                    'lines': lines,
                    'total_lines': len(lines),
                    'total_characters': len(content),
                    'encoding': encoding,
                    'extraction_time': datetime.now().isoformat()
                }
        except Exception as e:
            raise Exception(f"文字檔案解析失敗: {str(e)}")
    
    def _detect_encoding(self, file_path: str) -> str:
        """檢測檔案編碼"""
        encodings = ['utf-8', 'big5', 'gbk', 'gb2312', 'utf-16', 'ascii']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    file.read()
                return encoding
            except UnicodeDecodeError:
                continue
        
        return 'utf-8'  # 預設編碼 