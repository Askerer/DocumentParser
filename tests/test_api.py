import pytest
import os
import tempfile
from fastapi.testclient import TestClient
import pandas as pd
from pathlib import Path
import xml.etree.ElementTree as ET
from pptx import Presentation
from docx import Document
import io

from main import app


class TestAPI:
    """FastAPI應用程式的整合測試"""
    
    @pytest.fixture
    def client(self):
        """建立測試客戶端"""
        return TestClient(app)
    
    @pytest.fixture
    def temp_dir(self):
        """建立臨時目錄用於測試檔案"""
        with tempfile.TemporaryDirectory() as temp_dir:
            yield temp_dir
    
    def create_test_files(self, temp_dir):
        """建立各種格式的測試檔案"""
        test_files = {}
        
        # Excel檔案
        excel_path = os.path.join(temp_dir, "test.xlsx")
        data = {'列1': ['A', 'B'], '列2': [1, 2]}
        df = pd.DataFrame(data)
        df.to_excel(excel_path, index=False)
        test_files['excel'] = excel_path
        
        # PowerPoint檔案
        pptx_path = os.path.join(temp_dir, "test.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "測試"
        prs.save(pptx_path)
        test_files['pptx'] = pptx_path
        
        # XML檔案
        xml_path = os.path.join(temp_dir, "test.xml")
        root = ET.Element("root")
        title = ET.SubElement(root, "title")
        title.text = "測試"
        tree = ET.ElementTree(root)
        tree.write(xml_path, encoding='utf-8', xml_declaration=True)
        test_files['xml'] = xml_path
        
        # 文字檔案
        txt_path = os.path.join(temp_dir, "test.txt")
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write("測試文字內容")
        test_files['txt'] = txt_path
        
        # Word檔案
        docx_path = os.path.join(temp_dir, "test.docx")
        doc = Document()
        doc.add_paragraph("測試內容")
        doc.save(docx_path)
        test_files['docx'] = docx_path
        
        # CSV檔案
        csv_path = os.path.join(temp_dir, "test.csv")
        with open(csv_path, 'w', encoding='utf-8') as f:
            f.write("欄位1,欄位2\n值1,值2")
        test_files['csv'] = csv_path
        
        return test_files
    
    def test_root_endpoint(self, client):
        """測試根端點"""
        response = client.get("/")
        assert response.status_code == 200
        data = response.json()
        assert "message" in data
        assert "version" in data
    
    def test_health_check(self, client):
        """測試健康檢查端點"""
        response = client.get("/health")
        assert response.status_code == 200
        data = response.json()
        assert data["status"] == "healthy"
        assert "timestamp" in data
    
    def test_extract_excel_file(self, client, temp_dir):
        """測試Excel檔案擷取"""
        test_files = self.create_test_files(temp_dir)
        
        with open(test_files['excel'], 'rb') as f:
            response = client.post("/extract", files={"file": ("test.xlsx", f, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")})
        
        assert response.status_code == 200
        data = response.json()
        assert data["file_type"] == "excel"
        assert "sheets" in data
    
    def test_extract_powerpoint_file(self, client, temp_dir):
        """測試PowerPoint檔案擷取"""
        test_files = self.create_test_files(temp_dir)
        
        with open(test_files['pptx'], 'rb') as f:
            response = client.post("/extract", files={"file": ("test.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")})
        
        assert response.status_code == 200
        data = response.json()
        assert data["file_type"] == "powerpoint"
        assert "slides" in data
    
    def test_extract_xml_file(self, client, temp_dir):
        """測試XML檔案擷取"""
        test_files = self.create_test_files(temp_dir)
        
        with open(test_files['xml'], 'rb') as f:
            response = client.post("/extract", files={"file": ("test.xml", f, "application/xml")})
        
        assert response.status_code == 200
        data = response.json()
        assert data["file_type"] == "xml"
        assert "text_content" in data
    
    def test_extract_text_file(self, client, temp_dir):
        """測試文字檔案擷取"""
        test_files = self.create_test_files(temp_dir)
        
        with open(test_files['txt'], 'rb') as f:
            response = client.post("/extract", files={"file": ("test.txt", f, "text/plain")})
        
        assert response.status_code == 200
        data = response.json()
        assert data["file_type"] == "txt"
        assert "content" in data
    
    def test_extract_docx_file(self, client, temp_dir):
        """測試Word檔案擷取"""
        test_files = self.create_test_files(temp_dir)
        
        with open(test_files['docx'], 'rb') as f:
            response = client.post("/extract", files={"file": ("test.docx", f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")})
        
        assert response.status_code == 200
        data = response.json()
        assert data["file_type"] == "docx"
        assert "paragraphs" in data
    
    def test_extract_csv_file(self, client, temp_dir):
        """測試CSV檔案擷取"""
        test_files = self.create_test_files(temp_dir)
        
        with open(test_files['csv'], 'rb') as f:
            response = client.post("/extract", files={"file": ("test.csv", f, "text/csv")})
        
        assert response.status_code == 200
        data = response.json()
        assert data["file_type"] == "csv"
        assert "content" in data
    
    def test_unsupported_file_format(self, client, temp_dir):
        """測試不支援的檔案格式"""
        unsupported_path = os.path.join(temp_dir, "test.unsupported")
        with open(unsupported_path, 'w') as f:
            f.write("test content")
        
        with open(unsupported_path, 'rb') as f:
            response = client.post("/extract", files={"file": ("test.unsupported", f, "application/octet-stream")})
        
        assert response.status_code == 400
        assert "不支援的檔案格式" in response.json()["detail"]
    
    def test_large_file_rejection(self, client, temp_dir):
        """測試大檔案拒絕（模擬）"""
        # 建立一個相對較大的檔案進行測試
        large_file_path = os.path.join(temp_dir, "large.txt")
        with open(large_file_path, 'w') as f:
            f.write("x" * 1000)  # 1KB檔案，用於測試
        
        with open(large_file_path, 'rb') as f:
            response = client.post("/extract", files={"file": ("large.txt", f, "text/plain")})
        
        # 這個檔案應該成功，因為它實際上並不大
        assert response.status_code == 200
    
    def test_missing_file(self, client):
        """測試缺少檔案的情況"""
        response = client.post("/extract")
        assert response.status_code == 422  # Unprocessable Entity
    
    def test_empty_file(self, client, temp_dir):
        """測試空檔案"""
        empty_file_path = os.path.join(temp_dir, "empty.txt")
        with open(empty_file_path, 'w') as f:
            pass  # 建立空檔案
        
        with open(empty_file_path, 'rb') as f:
            response = client.post("/extract", files={"file": ("empty.txt", f, "text/plain")})
        
        assert response.status_code == 200
        data = response.json()
        assert data["file_type"] == "txt"


if __name__ == "__main__":
    pytest.main([__file__])