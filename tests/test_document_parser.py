import pytest
import os
import tempfile
import pandas as pd
from pathlib import Path
import xml.etree.ElementTree as ET
from pptx import Presentation
from docx import Document
import base64
from PIL import Image
import io

from document_parser import DocumentParser


class TestDocumentParser:
    """DocumentParser類別的單元測試"""
    
    @pytest.fixture
    def parser(self):
        """建立DocumentParser實例"""
        return DocumentParser()
    
    @pytest.fixture
    def temp_dir(self):
        """建立臨時目錄用於測試檔案"""
        with tempfile.TemporaryDirectory() as temp_dir:
            yield temp_dir
    
    def create_test_excel(self, file_path):
        """建立測試Excel檔案"""
        data = {
            '姓名': ['張三', '李四', '王五'],
            '年齡': [25, 30, 28],
            '城市': ['台北', '高雄', '台中']
        }
        df = pd.DataFrame(data)
        
        with pd.ExcelWriter(file_path, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='Sheet1', index=False)
            df.to_excel(writer, sheet_name='Sheet2', index=False)
    
    def create_test_powerpoint(self, file_path):
        """建立測試PowerPoint檔案"""
        prs = Presentation()
        
        # 第一張投影片
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "測試標題"
        slide1.placeholders[1].text = "測試內容"
        
        # 第二張投影片
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "第二張投影片"
        
        prs.save(file_path)
    
    def create_test_xml(self, file_path):
        """建立測試XML檔案"""
        root = ET.Element("root")
        
        title = ET.SubElement(root, "title")
        title.text = "測試XML檔案"
        
        description = ET.SubElement(root, "description") 
        description.text = "這是測試用的XML檔案"
        
        data = ET.SubElement(root, "data")
        item1 = ET.SubElement(data, "item", id="1")
        item1.text = "項目一"
        item2 = ET.SubElement(data, "item", id="2")
        item2.text = "項目二"
        
        tree = ET.ElementTree(root)
        tree.write(file_path, encoding='utf-8', xml_declaration=True)
    
    def create_test_docx(self, file_path):
        """建立測試Word檔案"""
        doc = Document()
        doc.add_heading('測試文件標題', 0)
        doc.add_paragraph('這是第一段測試內容。')
        doc.add_paragraph('這是第二段測試內容。')
        
        # 新增表格
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = '欄位1'
        table.cell(0, 1).text = '欄位2'
        table.cell(1, 0).text = '數據1'
        table.cell(1, 1).text = '數據2'
        
        doc.save(file_path)
    
    def create_test_txt(self, file_path):
        """建立測試文字檔案"""
        content = """這是測試文字檔案
包含多行內容
支援中文編碼
第四行內容"""
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content)
    
    def test_extract_excel_content(self, parser, temp_dir):
        """測試Excel檔案內容擷取"""
        file_path = os.path.join(temp_dir, "test.xlsx")
        self.create_test_excel(file_path)
        
        result = parser.extract_content(file_path, '.xlsx')
        
        assert result['file_type'] == 'excel'
        assert result['total_sheets'] == 2
        assert 'Sheet1' in result['sheet_names']
        assert 'Sheet2' in result['sheet_names']
        assert len(result['sheets']['Sheet1']['data']) == 3
        assert result['sheets']['Sheet1']['columns'] == ['姓名', '年齡', '城市']
    
    def test_extract_powerpoint_content(self, parser, temp_dir):
        """測試PowerPoint檔案內容擷取"""
        file_path = os.path.join(temp_dir, "test.pptx")
        self.create_test_powerpoint(file_path)
        
        result = parser.extract_content(file_path, '.pptx')
        
        assert result['file_type'] == 'powerpoint'
        assert result['total_slides'] == 2
        assert len(result['slides']) == 2
        assert any('測試標題' in str(slide) for slide in result['slides'])
    
    def test_extract_xml_content(self, parser, temp_dir):
        """測試XML檔案內容擷取"""
        file_path = os.path.join(temp_dir, "test.xml")
        self.create_test_xml(file_path)
        
        result = parser.extract_content(file_path, '.xml')
        
        assert result['file_type'] == 'xml'
        assert result['root_tag'] == 'root'
        assert result['total_elements'] > 0
        assert any('測試XML檔案' in item['content'] for item in result['text_content'])
    
    def test_extract_docx_content(self, parser, temp_dir):
        """測試Word檔案內容擷取"""
        file_path = os.path.join(temp_dir, "test.docx")
        self.create_test_docx(file_path)
        
        result = parser.extract_content(file_path, '.docx')
        
        assert result['file_type'] == 'docx'
        assert result['total_paragraphs'] > 0
        assert result['total_tables'] == 1
        assert any('測試文件標題' in para['text'] for para in result['paragraphs'])
    
    def test_extract_text_content(self, parser, temp_dir):
        """測試文字檔案內容擷取"""
        file_path = os.path.join(temp_dir, "test.txt")
        self.create_test_txt(file_path)
        
        result = parser.extract_content(file_path, '.txt')
        
        assert result['file_type'] == 'txt'
        assert result['total_lines'] == 4
        assert '這是測試文字檔案' in result['content']
        assert result['encoding'] == 'utf-8'
    
    def test_unsupported_file_format(self, parser, temp_dir):
        """測試不支援的檔案格式"""
        file_path = os.path.join(temp_dir, "test.unsupported")
        with open(file_path, 'w') as f:
            f.write("test")
        
        with pytest.raises(Exception) as exc_info:
            parser.extract_content(file_path, '.unsupported')
        
        assert "不支援的檔案格式" in str(exc_info.value)
    
    def test_detect_encoding(self, parser, temp_dir):
        """測試編碼檢測"""
        # 測試UTF-8編碼
        file_path_utf8 = os.path.join(temp_dir, "utf8.txt")
        with open(file_path_utf8, 'w', encoding='utf-8') as f:
            f.write("UTF-8 測試檔案")
        
        encoding = parser._detect_encoding(file_path_utf8)
        assert encoding == 'utf-8'
    
    def test_is_base64_image(self, parser):
        """測試Base64圖片檢測"""
        # 建立一個小的PNG圖片的base64編碼
        img = Image.new('RGB', (1, 1), color='red')
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        img_data = buffer.getvalue()
        base64_string = base64.b64encode(img_data).decode()
        
        # 測試非圖片文字
        assert not parser._is_base64_image("這不是圖片")
        assert not parser._is_base64_image("短文字")
        
        # 測試有效的base64但不是圖片
        non_image_b64 = base64.b64encode(b"not an image").decode()
        assert not parser._is_base64_image(non_image_b64)


if __name__ == "__main__":
    pytest.main([__file__])