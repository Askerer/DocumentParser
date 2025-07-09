#!/usr/bin/env python3
"""
文件內容擷取服務測試運行器

這個腳本提供多種測試模式：
1. 單元測試 - 測試DocumentParser類別
2. 整合測試 - 測試FastAPI應用程式
3. 效能測試 - 測試不同檔案大小和格式的效能
4. 自定義檔案測試 - 測試用戶指定的檔案
"""

import os
import sys
import argparse
import subprocess
import time
import json
from pathlib import Path
from typing import List, Dict, Any
import tempfile
import pandas as pd
import xml.etree.ElementTree as ET
from pptx import Presentation
from docx import Document

# 添加當前目錄到Python路徑
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from document_parser import DocumentParser


class TestRunner:
    """測試運行器類別"""
    
    def __init__(self):
        self.parser = DocumentParser()
        self.test_results = {}
        
    def run_unit_tests(self) -> Dict[str, Any]:
        """運行單元測試"""
        print("🧪 運行單元測試...")
        
        try:
            result = subprocess.run([
                sys.executable, "-m", "pytest", 
                "tests/test_document_parser.py", 
                "-v", "--tb=short"
            ], capture_output=True, text=True, cwd=".")
            
            return {
                "status": "success" if result.returncode == 0 else "failed",
                "returncode": result.returncode,
                "stdout": result.stdout,
                "stderr": result.stderr,
                "duration": 0
            }
        except Exception as e:
            return {
                "status": "error",
                "error": str(e),
                "duration": 0
            }
    
    def run_integration_tests(self) -> Dict[str, Any]:
        """運行整合測試"""
        print("🌐 運行整合測試...")
        
        try:
            result = subprocess.run([
                sys.executable, "-m", "pytest", 
                "tests/test_api.py", 
                "-v", "--tb=short"
            ], capture_output=True, text=True, cwd=".")
            
            return {
                "status": "success" if result.returncode == 0 else "failed",
                "returncode": result.returncode,
                "stdout": result.stdout,
                "stderr": result.stderr,
                "duration": 0
            }
        except Exception as e:
            return {
                "status": "error",
                "error": str(e),
                "duration": 0
            }
    
    def run_performance_tests(self) -> Dict[str, Any]:
        """運行效能測試"""
        print("⚡ 運行效能測試...")
        
        performance_results = {}
        
        with tempfile.TemporaryDirectory() as temp_dir:
            # 測試不同檔案格式的效能
            test_cases = [
                ("small_excel", self._create_small_excel),
                ("large_excel", self._create_large_excel),
                ("small_text", self._create_small_text),
                ("large_text", self._create_large_text),
                ("powerpoint", self._create_powerpoint),
                ("xml", self._create_xml),
                ("docx", self._create_docx)
            ]
            
            for test_name, create_func in test_cases:
                try:
                    file_path = create_func(temp_dir)
                    file_extension = Path(file_path).suffix
                    
                    # 測量處理時間
                    start_time = time.time()
                    result = self.parser.extract_content(file_path, file_extension)
                    end_time = time.time()
                    
                    file_size = os.path.getsize(file_path)
                    processing_time = end_time - start_time
                    
                    performance_results[test_name] = {
                        "file_size_bytes": file_size,
                        "file_size_mb": round(file_size / (1024 * 1024), 2),
                        "processing_time_seconds": round(processing_time, 3),
                        "throughput_mb_per_second": round((file_size / (1024 * 1024)) / processing_time, 2) if processing_time > 0 else 0,
                        "status": "success"
                    }
                    
                except Exception as e:
                    performance_results[test_name] = {
                        "status": "error",
                        "error": str(e)
                    }
        
        return {
            "status": "completed",
            "results": performance_results
        }
    
    def test_custom_files(self, file_paths: List[str]) -> Dict[str, Any]:
        """測試自定義檔案"""
        print(f"📁 測試自定義檔案: {', '.join(file_paths)}")
        
        custom_results = {}
        
        for file_path in file_paths:
            if not os.path.exists(file_path):
                custom_results[file_path] = {
                    "status": "error",
                    "error": f"檔案不存在: {file_path}"
                }
                continue
            
            try:
                file_extension = Path(file_path).suffix.lower()
                file_size = os.path.getsize(file_path)
                
                start_time = time.time()
                result = self.parser.extract_content(file_path, file_extension)
                end_time = time.time()
                
                processing_time = end_time - start_time
                
                custom_results[file_path] = {
                    "status": "success",
                    "file_size_bytes": file_size,
                    "file_size_mb": round(file_size / (1024 * 1024), 2),
                    "processing_time_seconds": round(processing_time, 3),
                    "file_type": result.get("file_type", "unknown"),
                    "content_summary": self._summarize_content(result)
                }
                
            except Exception as e:
                custom_results[file_path] = {
                    "status": "error",
                    "error": str(e)
                }
        
        return {
            "status": "completed",
            "results": custom_results
        }
    
    def run_all_tests(self) -> Dict[str, Any]:
        """運行所有測試"""
        print("🚀 運行完整測試套件...")
        
        all_results = {}
        
        # 運行單元測試
        all_results["unit_tests"] = self.run_unit_tests()
        
        # 運行整合測試
        all_results["integration_tests"] = self.run_integration_tests()
        
        # 運行效能測試
        all_results["performance_tests"] = self.run_performance_tests()
        
        return all_results
    
    def generate_report(self, results: Dict[str, Any], output_file: str = None):
        """生成測試報告"""
        if output_file:
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(results, f, ensure_ascii=False, indent=2)
            print(f"📊 測試報告已儲存至: {output_file}")
        
        # 控制台報告
        print("\n" + "="*80)
        print("📈 測試結果摘要")
        print("="*80)
        
        for test_type, result in results.items():
            print(f"\n{test_type.upper()}:")
            if isinstance(result, dict) and "status" in result:
                print(f"  狀態: {result['status']}")
                if result["status"] == "failed" and "stderr" in result:
                    print(f"  錯誤: {result['stderr'][:200]}...")
                elif result["status"] == "completed" and "results" in result:
                    print(f"  測試項目數: {len(result['results'])}")
                    success_count = sum(1 for r in result['results'].values() if r.get('status') == 'success')
                    print(f"  成功: {success_count}/{len(result['results'])}")
    
    # 測試檔案建立方法
    def _create_small_excel(self, temp_dir):
        """建立小型Excel檔案"""
        file_path = os.path.join(temp_dir, "small.xlsx")
        data = {'列1': list(range(10)), '列2': [f'值{i}' for i in range(10)]}
        df = pd.DataFrame(data)
        df.to_excel(file_path, index=False)
        return file_path
    
    def _create_large_excel(self, temp_dir):
        """建立大型Excel檔案"""
        file_path = os.path.join(temp_dir, "large.xlsx")
        data = {'列1': list(range(1000)), '列2': [f'值{i}' for i in range(1000)]}
        df = pd.DataFrame(data)
        df.to_excel(file_path, index=False)
        return file_path
    
    def _create_small_text(self, temp_dir):
        """建立小型文字檔案"""
        file_path = os.path.join(temp_dir, "small.txt")
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write("小型測試檔案\n" * 10)
        return file_path
    
    def _create_large_text(self, temp_dir):
        """建立大型文字檔案"""
        file_path = os.path.join(temp_dir, "large.txt")
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write("大型測試檔案內容\n" * 1000)
        return file_path
    
    def _create_powerpoint(self, temp_dir):
        """建立PowerPoint檔案"""
        file_path = os.path.join(temp_dir, "test.pptx")
        prs = Presentation()
        for i in range(5):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = f"投影片 {i+1}"
        prs.save(file_path)
        return file_path
    
    def _create_xml(self, temp_dir):
        """建立XML檔案"""
        file_path = os.path.join(temp_dir, "test.xml")
        root = ET.Element("root")
        for i in range(10):
            item = ET.SubElement(root, "item", id=str(i))
            item.text = f"項目 {i}"
        tree = ET.ElementTree(root)
        tree.write(file_path, encoding='utf-8', xml_declaration=True)
        return file_path
    
    def _create_docx(self, temp_dir):
        """建立Word檔案"""
        file_path = os.path.join(temp_dir, "test.docx")
        doc = Document()
        doc.add_heading('測試文件', 0)
        for i in range(10):
            doc.add_paragraph(f'段落 {i+1}: 這是測試內容。')
        doc.save(file_path)
        return file_path
    
    def _summarize_content(self, result):
        """總結內容擷取結果"""
        summary = {"type": result.get("file_type", "unknown")}
        
        if "sheets" in result:
            summary["sheets_count"] = result.get("total_sheets", 0)
        elif "slides" in result:
            summary["slides_count"] = result.get("total_slides", 0)
        elif "paragraphs" in result:
            summary["paragraphs_count"] = result.get("total_paragraphs", 0)
        elif "content" in result:
            summary["characters_count"] = len(result["content"])
        
        return summary


def main():
    """主程式"""
    parser = argparse.ArgumentParser(description="文件內容擷取服務測試運行器")
    parser.add_argument("--mode", choices=["unit", "integration", "performance", "custom", "all"], 
                       default="all", help="測試模式")
    parser.add_argument("--files", nargs="+", help="自定義測試檔案路徑（僅用於custom模式）")
    parser.add_argument("--output", help="輸出報告檔案路徑")
    
    args = parser.parse_args()
    
    runner = TestRunner()
    
    if args.mode == "unit":
        results = {"unit_tests": runner.run_unit_tests()}
    elif args.mode == "integration":
        results = {"integration_tests": runner.run_integration_tests()}
    elif args.mode == "performance":
        results = {"performance_tests": runner.run_performance_tests()}
    elif args.mode == "custom":
        if not args.files:
            print("❌ 自定義模式需要指定檔案路徑")
            sys.exit(1)
        results = {"custom_tests": runner.test_custom_files(args.files)}
    else:  # all
        results = runner.run_all_tests()
    
    runner.generate_report(results, args.output)


if __name__ == "__main__":
    main()